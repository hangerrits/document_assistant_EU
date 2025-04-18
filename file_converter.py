#!/usr/bin/env python3
"""
File Converter Module

This module provides functions to convert various file types to PDF format.
It supports common document formats like DOCX, PPTX, TXT, etc.
"""

import os
import logging
import tempfile
import shutil
import re
from typing import Optional, Tuple

# Import required libraries for file conversion
import pypdfium2 as pdfium
from docx import Document
from pptx import Presentation
from PIL import Image
import io

# Import ReportLab for better PDF generation with Unicode support
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

def get_file_extension(file_path: str) -> str:
    """Get the lowercase extension of a file without the dot."""
    _, ext = os.path.splitext(file_path)
    return ext.lower().lstrip('.')

def is_pdf_file(file_path: str) -> bool:
    """Check if a file is a PDF based on its extension."""
    return get_file_extension(file_path) == 'pdf'

def get_base_filename_no_ext(file_path: str) -> str:
    """Get the base filename without the extension."""
    base_name = os.path.basename(file_path)
    name, _ = os.path.splitext(base_name)
    return name

def clean_filename_for_title(filename_part: str) -> str:
    """
    Clean a filename part (without extension) to make it suitable for display as a document title.

    Args:
        filename_part: Original filename part without extension.

    Returns:
        Cleaned name for title.
    """
    name = filename_part

    # Remove temp_ prefix if present (from the upload process)
    if name.startswith('temp_'):
        name = name[5:]

    # Replace underscores and hyphens with spaces
    name = name.replace('_', ' ').replace('-', ' ')

    # Clean up multiple spaces
    name = re.sub(r'\s+', ' ', name).strip()

    # Handle cases where cleaning results in an empty string
    if not name:
        return "Untitled Document" # Or some other default

    return name

def get_document_title(file_path: str) -> str:
    """Extract a clean document title from the file path."""
    base_name_no_ext = get_base_filename_no_ext(file_path)
    return clean_filename_for_title(base_name_no_ext)

def convert_to_pdf(input_file: str) -> Tuple[bool, Optional[str]]:
    """
    Convert a file to PDF format based on its extension.

    Args:
        input_file: Path to the input file

    Returns:
        Tuple containing:
            - Success status (bool)
            - Path to the converted PDF file (str) or None if conversion failed
    """
    if not os.path.exists(input_file):
        logger.error(f"Input file does not exist: {input_file}")
        return False, None

    # If file is already a PDF, return it as is
    if is_pdf_file(input_file):
        logger.info(f"File is already a PDF, no conversion needed: {input_file}")
        # Even if it's already PDF, we might want to ensure it's in a temp dir
        # for consistency, but for now, let's just return it.
        # Consider copying it to a temp dir if downstream processing expects it there.
        return True, input_file

    # Get file extension to determine conversion method
    ext = get_file_extension(input_file)

    # Create a temporary directory for conversion
    temp_dir = tempfile.mkdtemp()

    # --- CHANGE: Use original base name for the output file path ---
    original_base_name_no_ext = get_base_filename_no_ext(input_file)
    if not original_base_name_no_ext: # Handle cases like '.docx'
        original_base_name_no_ext = "untitled"
    output_filename = f"{original_base_name_no_ext}.pdf"
    output_file = os.path.join(temp_dir, output_filename)
    # --- END CHANGE ---

    # Get document title for metadata (use the cleaning function here)
    document_title = get_document_title(input_file)

    try:
        success = False
        # Choose conversion method based on file extension
        if ext == 'docx':
            success = convert_docx_to_pdf(input_file, output_file, document_title)
        elif ext == 'pptx':
            success = convert_pptx_to_pdf(input_file, output_file, document_title)
        elif ext in ['txt', 'md', 'csv']:
            success = convert_text_to_pdf(input_file, output_file, document_title)
        elif ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'tif', 'webp']:
            success = convert_image_to_pdf(input_file, output_file, document_title)
        else:
            logger.warning(f"Unsupported file type for conversion: {ext}")
            shutil.rmtree(temp_dir) # Clean up temp dir on failure
            return False, None

        if success:
            logger.info(f"Successfully converted {input_file} to PDF: {output_file}")
            return True, output_file
        else:
            logger.error(f"Failed to convert {input_file} to PDF")
            shutil.rmtree(temp_dir) # Clean up temp dir on failure
            return False, None

    except Exception as e:
        logger.error(f"Error converting {input_file} to PDF: {e}", exc_info=True)
        if os.path.exists(temp_dir): # Ensure cleanup even on exception
             shutil.rmtree(temp_dir)
        return False, None

# --- Functions convert_docx_to_pdf, convert_pptx_to_pdf, convert_text_to_pdf, convert_image_to_pdf remain largely the same ---
# Ensure they correctly use the `document_title` parameter for PDF metadata Title
# and write to the exact `output_file` path provided.

def convert_docx_to_pdf(input_file: str, output_file: str, document_title: str) -> bool:
    """
    Convert a DOCX file to PDF using ReportLab.
    Uses the `document_title` for PDF metadata title.
    Writes the output to `output_file`.
    """
    try:
        doc = Document(input_file)
        pdf_doc = SimpleDocTemplate(
            output_file, # Use the exact output path
            pagesize=letter,
            rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72,
            title=document_title, # Use cleaned title for metadata
            author="Document Assistant",
            subject=f"Converted from {os.path.basename(input_file)}",
            creator="PDF Conversion Tool"
        )
        styles = getSampleStyleSheet()
        normal_style = styles['Normal']
        heading1_style = styles['Heading1']
        heading2_style = styles['Heading2']
        flowables = []
        flowables.append(Paragraph(document_title, heading1_style)) # Use cleaned title maybe as H1
        flowables.append(Spacer(1, 0.3 * inch))
        for paragraph in doc.paragraphs:
            if not paragraph.text.strip():
                flowables.append(Spacer(1, 0.2 * inch))
                continue
            style = normal_style
            if paragraph.style.name.startswith('Heading1') or paragraph.style.name.startswith('Title'):
                style = heading1_style
            elif paragraph.style.name.startswith('Heading2'):
                style = heading2_style
            # Basic handling, might need more sophisticated style mapping
            # Consider escaping HTML-like characters in paragraph.text if needed
            flowables.append(Paragraph(paragraph.text.replace('&', '&').replace('<', '<').replace('>', '>'), style))
            flowables.append(Spacer(1, 0.1 * inch))
        pdf_doc.build(flowables)
        return os.path.exists(output_file) and os.path.getsize(output_file) > 0
    except Exception as e:
        logger.error(f"Error converting DOCX {input_file} to PDF: {e}", exc_info=True)
        return False

def convert_pptx_to_pdf(input_file: str, output_file: str, document_title: str) -> bool:
    """
    Convert a PPTX file to PDF.
    Uses the `document_title` for PDF metadata title.
    Writes the output to `output_file`.
    """
    try:
        prs = Presentation(input_file)
        pdf_doc = SimpleDocTemplate(
            output_file, # Use the exact output path
            pagesize=letter, # Consider landscape? prs.slide_width/height?
            rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36,
            title=document_title, # Use cleaned title for metadata
            author="Document Assistant",
            subject=f"Converted from {os.path.basename(input_file)}",
            creator="PDF Conversion Tool"
        )
        styles = getSampleStyleSheet()
        title_style = styles['Title']
        normal_style = styles['Normal']
        flowables = []
        flowables.append(Paragraph(document_title, title_style)) # Use cleaned title
        flowables.append(Spacer(1, 0.3 * inch))
        for i, slide in enumerate(prs.slides):
            slide_title = f"Slide {i+1}"
            if slide.shapes.title:
                slide_title = slide.shapes.title.text
            flowables.append(Paragraph(slide_title, styles['Heading2'])) # Use a heading for slide title
            flowables.append(Spacer(1, 0.2 * inch))
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                     # Consider escaping HTML-like characters
                    flowables.append(Paragraph(shape.text.replace('&', '&').replace('<', '<').replace('>', '>'), normal_style))
                    flowables.append(Spacer(1, 0.1 * inch))
            # Add page break (or similar separator) after each slide's content
            # Using Spacer might not force a page break effectively,
            # Consider using PageBreak from reportlab.platypus if needed
            flowables.append(Spacer(1, 0.5 * inch)) # Visual separator
            # from reportlab.platypus import PageBreak
            # if i < len(prs.slides) - 1:
            #     flowables.append(PageBreak())

        pdf_doc.build(flowables)
        return os.path.exists(output_file) and os.path.getsize(output_file) > 0
    except Exception as e:
        logger.error(f"Error converting PPTX {input_file} to PDF: {e}", exc_info=True)
        return False

def convert_text_to_pdf(input_file: str, output_file: str, document_title: str) -> bool:
    """
    Convert a text file (TXT, MD, CSV) to PDF using ReportLab.
    Uses the `document_title` for PDF metadata title.
    Writes the output to `output_file`.
    """
    try:
        with open(input_file, 'r', encoding='utf-8', errors='replace') as f: # Added errors='replace'
            text = f.read()
        pdf_doc = SimpleDocTemplate(
            output_file, # Use the exact output path
            pagesize=letter,
            rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72,
            title=document_title, # Use cleaned title for metadata
            author="Document Assistant",
            subject=f"Converted from {os.path.basename(input_file)}",
            creator="PDF Conversion Tool"
        )
        styles = getSampleStyleSheet()
        # Use a monospaced font for text files, potentially?
        # pdfmetrics.registerFont(TTFont('DejaVuSansMono', 'DejaVuSansMono.ttf')) # Need font file
        # mono_style = ParagraphStyle(name='Mono', parent=styles['Normal'], fontName='DejaVuSansMono')
        normal_style = styles['Code'] # Use Code style for better preformatted look
        title_style = styles['Title']
        flowables = []
        flowables.append(Paragraph(document_title, title_style)) # Use cleaned title
        flowables.append(Spacer(1, 0.3 * inch))

        # Use Paragraph with preformatted=True or specific style
        # Need to escape HTML entities
        escaped_text = text.replace('&', '&').replace('<', '<').replace('>', '>')
        # Split into paragraphs based on double newline, or just treat line by line
        # Using Paragraph preserves line breaks within the paragraph text itself
        for line in escaped_text.splitlines():
             flowables.append(Paragraph(line, normal_style))
             # flowables.append(Spacer(1, 0.05 * inch)) # Smaller spacer between lines if needed


        # Alternative: Use Preformatted style if available or define one
        # from reportlab.platypus import Preformatted
        # flowables.append(Preformatted(escaped_text, normal_style))


        pdf_doc.build(flowables)
        return os.path.exists(output_file) and os.path.getsize(output_file) > 0
    except Exception as e:
        logger.error(f"Error converting text {input_file} to PDF: {e}", exc_info=True)
        return False

def convert_image_to_pdf(input_file: str, output_file: str, document_title: str) -> bool:
    """
    Convert an image file to PDF with proper metadata.
    Uses the `document_title` for PDF metadata title.
    Writes the output to `output_file`.
    """
    temp_img_path = None # Initialize
    try:
        img = Image.open(input_file)
        img_width, img_height = img.size

        # Determine page size based on image, or fit to standard page
        page_width, page_height = letter # Use standard letter size
        max_content_width = page_width - 72 # Margins (36*2)
        max_content_height = page_height - 72 # Margins (36*2)

        scale_w = max_content_width / img_width
        scale_h = max_content_height / img_height
        scale = min(scale_w, scale_h, 1.0) # Don't scale up, max scale is 1

        draw_width = img_width * scale
        draw_height = img_height * scale

        # Convert RGBA to RGB if needed for ReportLab Image
        # ReportLab's RLImage handles transparency better now, but saving intermediary might need RGB
        temp_img_file = None
        if img.mode == 'RGBA':
            # Try preserving transparency if possible, otherwise convert
            # RLImage might handle PIL Image object directly
            pass # Let RLImage handle PIL object directly if possible
        elif img.mode == 'P': # Palette mode, convert to RGB
             img = img.convert('RGB')

        # Create PDF with ReportLab
        pdf_doc = SimpleDocTemplate(
            output_file, # Use the exact output path
            pagesize=(page_width, page_height),
            rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36,
            title=document_title, # Use cleaned title for metadata
            author="Document Assistant",
            subject=f"Converted from {os.path.basename(input_file)}",
            creator="PDF Conversion Tool"
        )
        styles = getSampleStyleSheet()
        title_style = styles['Title']
        flowables = []
        flowables.append(Paragraph(document_title, title_style)) # Use cleaned title
        flowables.append(Spacer(1, 0.3 * inch))

        # Use RLImage with the PIL image object directly if possible
        rl_img = RLImage(img, width=draw_width, height=draw_height)
        # If RLImage fails with PIL object, save to temp file first
        # try:
        #    rl_img = RLImage(img, width=draw_width, height=draw_height)
        # except:
        #    logger.warning("RLImage failed with PIL object, saving temp image.")
        #    temp_img_format = "JPEG" if img.mode != 'RGBA' else "PNG"
        #    temp_img_path = output_file + f".temp.{temp_img_format.lower()}"
        #    img.save(temp_img_path, temp_img_format)
        #    rl_img = RLImage(temp_img_path, width=draw_width, height=draw_height)

        flowables.append(rl_img)
        pdf_doc.build(flowables)

        # Clean up temporary image file if created
        if temp_img_path and os.path.exists(temp_img_path):
            os.remove(temp_img_path)

        return os.path.exists(output_file) and os.path.getsize(output_file) > 0
    except Exception as e:
        logger.error(f"Error converting image {input_file} to PDF: {e}", exc_info=True)
        if temp_img_path and os.path.exists(temp_img_path): # Cleanup on error too
            os.remove(temp_img_path)
        return False
    finally:
        if 'img' in locals() and hasattr(img, 'close'):
             img.close() # Close the PIL image object


def convert_file_for_upload(input_file: str) -> Tuple[bool, Optional[str], Optional[str]]:
    """
    Convert a file to PDF if needed and prepare it for upload.

    Args:
        input_file: Path to the input file

    Returns:
        Tuple containing:
            - Success status (bool)
            - Path to the file to be uploaded (str) - either original PDF or converted PDF in a temp location
            - Suggested final filename (str) - preserving original base name + .pdf, or None if conversion failed
    """
    original_basename = os.path.basename(input_file)

    # If file is already a PDF, return its path and original name
    if is_pdf_file(input_file):
        # It might be beneficial to copy even existing PDFs to a single
        # managed temp location before upload, but for now, we return original.
        logger.info(f"File {original_basename} is already PDF. No conversion needed.")
        return True, input_file, original_basename

    # --- CHANGE: Determine the desired *final* filename based on original base name ---
    original_base_name_no_ext = get_base_filename_no_ext(input_file)
    if not original_base_name_no_ext: # Handle cases like '.docx'
         original_base_name_no_ext = "untitled"
    suggested_final_filename = f"{original_base_name_no_ext}.pdf"
    # --- END CHANGE ---

    # Convert file to PDF (will be placed in a temporary directory)
    success, pdf_path = convert_to_pdf(input_file)

    if success and pdf_path:
        # pdf_path is the actual location of the converted file (e.g., /tmp/xyz/test.pdf)
        # suggested_final_filename is the name the file should have after upload (e.g., test.pdf)
        return True, pdf_path, suggested_final_filename
    else:
        logger.error(f"Failed to convert {original_basename} for upload.")
        return False, None, None

# For testing purposes
if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python file_converter.py <input_file>")
        sys.exit(1)

    input_file_arg = sys.argv[1]

    # Create dummy files for testing if they don't exist
    test_files_to_create = {
        "test.docx": "This is a test document.",
        "presentation test.pptx": None, # Requires python-pptx to create simply
        "data_file.txt": "Line 1\nLine 2\n---\nEnd of data.",
        "image-test.jpg": None, # Requires Pillow to create simply
        "already.pdf": None # Assume this exists or create manually
    }

    if not os.path.exists(input_file_arg):
        print(f"Input file {input_file_arg} not found. Attempting to create dummy files...")
        try:
            # Dummy DOCX
            if not os.path.exists("test.docx"):
                doc = Document()
                doc.add_paragraph(test_files_to_create["test.docx"])
                doc.save("test.docx")
                print("Created dummy test.docx")

            # Dummy PPTX (basic)
            if not os.path.exists("presentation test.pptx"):
                 try:
                     prs = Presentation()
                     slide = prs.slides.add_slide(prs.slide_layouts[0])
                     title = slide.shapes.title
                     title.text = "Test Presentation Slide"
                     prs.save("presentation test.pptx")
                     print("Created dummy presentation test.pptx")
                 except Exception as e:
                     print(f"Could not create dummy pptx (python-pptx needed?): {e}")


            # Dummy TXT
            if not os.path.exists("data_file.txt"):
                with open("data_file.txt", "w") as f:
                    f.write(test_files_to_create["data_file.txt"])
                print("Created dummy data_file.txt")

            # Dummy JPG (basic red square)
            if not os.path.exists("image-test.jpg"):
                try:
                    img = Image.new('RGB', (60, 30), color = 'red')
                    img.save('image-test.jpg')
                    print("Created dummy image-test.jpg")
                except Exception as e:
                     print(f"Could not create dummy jpg (Pillow needed?): {e}")

            # Dummy PDF (very basic)
            if not os.path.exists("already.pdf"):
                 try:
                     pdf_doc = SimpleDocTemplate("already.pdf", pagesize=letter)
                     flowables = [Paragraph("This is already a PDF.")]
                     pdf_doc.build(flowables)
                     print("Created dummy already.pdf")
                 except Exception as e:
                     print(f"Could not create dummy pdf (reportlab needed?): {e}")


        except Exception as e:
            print(f"Error creating dummy files: {e}")
            print(f"Please ensure {input_file_arg} exists.")
            # sys.exit(1) # Continue execution if the specified file exists anyway

        # Recheck if the specific input file exists now
        if not os.path.exists(input_file_arg):
             print(f"Input file {input_file_arg} still not found after attempting creation.")
             sys.exit(1)


    print(f"\n--- Testing conversion for: {input_file_arg} ---")
    success_upload, file_to_upload_path, suggested_name = convert_file_for_upload(input_file_arg)

    if success_upload:
        print(f"\nConversion/Preparation successful!")
        print(f"  File ready for processing: {file_to_upload_path}")
        print(f"  Suggested final filename:  {suggested_name}")

        # Optional: verify the created/returned file exists
        if os.path.exists(file_to_upload_path):
             print(f"  Verified file exists at path.")
             # If it's a PDF, try opening with pdfium to check basic validity
             if is_pdf_file(file_to_upload_path):
                 try:
                     pdf = pdfium.PdfDocument(file_to_upload_path)
                     print(f"  PDF check: Found {len(pdf)} pages.")
                     pdf.close()
                 except Exception as pdf_err:
                     print(f"  PDF check: Error opening PDF - {pdf_err}")
        else:
             print(f"  ERROR: File path returned does not exist!")

        # Clean up the temporary file/directory if it's not the original input
        if file_to_upload_path != input_file_arg and os.path.exists(file_to_upload_path):
            temp_dir_path = os.path.dirname(file_to_upload_path)
            # Be careful with recursive delete! Ensure it's really a temp dir.
            # A safer approach is to register temp dirs and clean them up later.
            # For this example, let's assume the path structure is from mkdtemp.
            if tempfile.gettempdir() in temp_dir_path:
                 print(f"  Cleaning up temporary file/directory: {temp_dir_path}")
                 try:
                     # If pdf_path is the file, delete the containing directory
                     shutil.rmtree(temp_dir_path)
                 except Exception as clean_err:
                     print(f"  Error during cleanup: {clean_err}")
            else:
                print(f"  Skipping cleanup - path doesn't look like a standard temp dir: {file_to_upload_path}")


    else:
        print("\nConversion/Preparation failed.")